from pathlib import Path
from typing import List, Union, Dict, Any
import logging
from PIL import Image
import fitz  # PyMuPDF
import tempfile
import subprocess
import os
from io import BytesIO
from pptx import Presentation
import re

logger = logging.getLogger(__name__)


class PPTParser:
    """PPT/PPTX parser that converts PowerPoint files to images"""
    
    def __init__(self, dpi: int = 200, use_native_pptx: bool = True, 
                 preserve_fonts: bool = True, antialiasing: bool = True):
        """
        Initialize PPT Parser
        
        Args:
            dpi: DPI for image conversion (default: 200)
            use_native_pptx: Use python-pptx for PPTX files when possible
            preserve_fonts: Try to preserve original fonts
            antialiasing: Enable antialiasing for better quality
        """
        self.dpi = dpi
        self.use_native_pptx = use_native_pptx
        self.preserve_fonts = preserve_fonts
        self.antialiasing = antialiasing
        self._check_dependencies()
        
    def parse(
        self, 
        ppt_path: Union[str, Path], 
        output_dir: Union[str, Path],
        image_format: str = "png",
        prefix: str = "slide"
    ) -> List[Path]:
        """
        Parse PPT/PPTX file and convert each slide to an image
        
        Args:
            ppt_path: Path to PPT/PPTX file
            output_dir: Directory to save images
            image_format: Image format (png, jpg, jpeg)
            prefix: Prefix for image filenames
            
        Returns:
            List of paths to generated images
        """
        ppt_path = Path(ppt_path)
        output_dir = Path(output_dir)
        
        if not ppt_path.exists():
            raise FileNotFoundError(f"PPT file not found: {ppt_path}")
        
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Check file extension
        if ppt_path.suffix.lower() == '.pptx':
            return self._parse_pptx(ppt_path, output_dir, image_format, prefix)
        elif ppt_path.suffix.lower() == '.ppt':
            return self._parse_ppt(ppt_path, output_dir, image_format, prefix)
        else:
            raise ValueError(f"Unsupported file format: {ppt_path.suffix}")
            
    def _parse_pptx(
        self,
        pptx_path: Path,
        output_dir: Path,
        image_format: str,
        prefix: str
    ) -> List[Path]:
        """Parse PPTX file with enhanced options"""
        try:
            # Try to extract slide information using python-pptx first
            if self.use_native_pptx:
                try:
                    prs = Presentation(str(pptx_path))
                    slide_count = len(prs.slides)
                    logger.info(f"PPTX file has {slide_count} slides")
                    
                    # Extract slide dimensions for optimal rendering
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    aspect_ratio = slide_width / slide_height
                    logger.info(f"Slide dimensions: {slide_width}x{slide_height}, aspect ratio: {aspect_ratio:.2f}")
                    
                    # Check for potential text rendering issues
                    has_number_patterns = False
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                # Check for patterns like "01", "02", etc.
                                if re.search(r'\b0[1-9]\b', shape.text):
                                    has_number_patterns = True
                                    logger.warning(
                                        f"⚠️  Detected number patterns (01, 02, etc.) that may render with spaces. "
                                        f"Consider using Microsoft PowerPoint for conversion if text accuracy is critical."
                                    )
                                    break
                        if has_number_patterns:
                            break
                            
                except Exception as e:
                    logger.warning(f"Could not parse PPTX with python-pptx: {e}")
            
            # Convert PPTX to PDF using enhanced LibreOffice options
            # Create temp file path without using NamedTemporaryFile to avoid premature deletion
            temp_pdf_fd, temp_pdf_path = tempfile.mkstemp(suffix='.pdf')
            os.close(temp_pdf_fd)  # Close file descriptor immediately
            pdf_path = Path(temp_pdf_path)
            
            try:
                # Convert PPTX to PDF with enhanced quality
                self._convert_to_pdf(pptx_path, pdf_path)
                
                # Convert PDF to images with enhanced settings
                images = self._pdf_to_images(pdf_path, output_dir, image_format, prefix)
                
            finally:
                # Clean up temporary PDF
                if pdf_path.exists():
                    try:
                        pdf_path.unlink()
                    except Exception as e:
                        logger.warning(f"Could not delete temporary PDF: {e}")
            
            return images
            
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            raise
            
    def _parse_ppt(
        self,
        ppt_path: Path,
        output_dir: Path,
        image_format: str,
        prefix: str
    ) -> List[Path]:
        """Parse PPT file by converting to PDF first"""
        try:
            # Create temp file path without using NamedTemporaryFile
            temp_pdf_fd, temp_pdf_path = tempfile.mkstemp(suffix='.pdf')
            os.close(temp_pdf_fd)  # Close file descriptor immediately
            pdf_path = Path(temp_pdf_path)
            
            try:
                self._convert_to_pdf(ppt_path, pdf_path)
                
                # Convert PDF to images
                images = self._pdf_to_images(pdf_path, output_dir, image_format, prefix)
                
            finally:
                # Clean up temporary PDF
                if pdf_path.exists():
                    try:
                        pdf_path.unlink()
                    except Exception as e:
                        logger.warning(f"Could not delete temporary PDF: {e}")
            
            return images
            
        except Exception as e:
            logger.error(f"Error parsing PPT: {e}")
            raise
            
    def _extract_slide_texts(self, pptx_path: Path) -> Dict[int, str]:
        """Extract text content from each slide for text preservation"""
        slide_texts = {}
        
        try:
            prs = Presentation(str(pptx_path))
            for slide_idx, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Preserve original text without modifications
                        texts.append(shape.text.strip())
                slide_texts[slide_idx] = "\n".join(texts)
                
            logger.info(f"Extracted text from {len(slide_texts)} slides")
            return slide_texts
            
        except Exception as e:
            logger.warning(f"Could not extract texts from PPTX: {e}")
            return {}

    def _convert_to_pdf(self, input_path: Path, output_path: Path):
        """Convert PPT/PPTX to PDF using LibreOffice with enhanced options"""
        
        # First, check if input file exists and is not empty
        if not input_path.exists():
            raise FileNotFoundError(f"Input file not found: {input_path}")
        
        if input_path.stat().st_size == 0:
            raise ValueError(f"Input file is empty: {input_path}")
        
        # Try simpler conversion first
        success = False
        
        # Method 1: Simple LibreOffice conversion
        if self.has_libreoffice:
            cmd = [
                'libreoffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', str(output_path.parent),
                str(input_path)
            ]
            
            try:
                logger.info(f"Converting {input_path} to PDF using LibreOffice...")
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
                
                # Check if conversion was successful
                expected_pdf = output_path.parent / f"{input_path.stem}.pdf"
                if expected_pdf.exists() and expected_pdf.stat().st_size > 0:
                    if expected_pdf != output_path:
                        expected_pdf.rename(output_path)
                    success = True
                    logger.info("PDF conversion successful")
                else:
                    logger.warning(f"LibreOffice conversion failed or produced empty file. stderr: {result.stderr}")
                    
            except subprocess.TimeoutExpired:
                logger.warning("LibreOffice conversion timed out")
            except Exception as e:
                logger.warning(f"LibreOffice conversion error: {e}")
        
        # Method 2: Try unoconv if available and first method failed
        if not success and self.has_unoconv:
            # unoconv with better font handling
            cmd = [
                'unoconv',
                '-f', 'pdf',
                '-e', 'SelectPdfVersion=1',  # PDF 1.5
                '-e', 'EmbedStandardFonts=true',  # Embed fonts
                '-e', 'UseTaggedPDF=true',  # Better text structure
                '-o', str(output_path),
                str(input_path)
            ]
            
            try:
                logger.info("Trying unoconv with enhanced options...")
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=60, check=True)
                
                if output_path.exists() and output_path.stat().st_size > 0:
                    success = True
                    logger.info("PDF conversion successful with unoconv")
                    
            except Exception as e:
                logger.warning(f"Unoconv conversion error: {e}")
        
        # Method 3: Try with enhanced options if basic conversion failed
        if not success and self.has_libreoffice:
            # Clean up any existing empty PDF
            if output_path.exists():
                output_path.unlink()
                
            # Try with more compatible options
            cmd = [
                'libreoffice',
                '--headless',
                '--invisible',
                '--nodefault',
                '--nofirststartwizard',
                '--nolockcheck',
                '--nologo',
                '--norestore',
                '--convert-to', 'pdf:writer_pdf_Export',
                '--outdir', str(output_path.parent),
                str(input_path)
            ]
            
            try:
                logger.info("Trying LibreOffice with compatibility options...")
                env = subprocess.os.environ.copy()
                # Add HOME environment variable if not set (common issue in containers)
                if 'HOME' not in env:
                    env['HOME'] = '/tmp'
                    
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=90, env=env)
                
                expected_pdf = output_path.parent / f"{input_path.stem}.pdf"
                if expected_pdf.exists() and expected_pdf.stat().st_size > 0:
                    if expected_pdf != output_path:
                        expected_pdf.rename(output_path)
                    success = True
                    logger.info("PDF conversion successful with compatibility options")
                else:
                    logger.error(f"All conversion attempts failed. stderr: {result.stderr}")
                    
            except Exception as e:
                logger.error(f"Final conversion attempt error: {e}")
        
        # Validate the output
        if not success or not output_path.exists() or output_path.stat().st_size == 0:
            error_msg = "Failed to convert PPT to PDF. "
            if not self.has_libreoffice and not self.has_unoconv:
                error_msg += "Please install LibreOffice: sudo apt-get install libreoffice"
            else:
                error_msg += "The file may be corrupted or in an unsupported format."
            
            # Clean up empty file if it exists
            if output_path.exists() and output_path.stat().st_size == 0:
                output_path.unlink()
                
            raise Exception(error_msg)
    
    def _check_dependencies(self):
        """Check if required dependencies are available"""
        try:
            subprocess.run(['libreoffice', '--version'], capture_output=True, check=True)
            self.has_libreoffice = True
        except (subprocess.CalledProcessError, FileNotFoundError):
            self.has_libreoffice = False
            logger.warning("LibreOffice not found. PPT conversion quality may be limited.")
            
        # Check for unoconv as alternative
        try:
            subprocess.run(['unoconv', '--version'], capture_output=True, check=True)
            self.has_unoconv = True
        except (subprocess.CalledProcessError, FileNotFoundError):
            self.has_unoconv = False
            
    def _pdf_to_images(
        self,
        pdf_path: Path,
        output_dir: Path,
        image_format: str,
        prefix: str
    ) -> List[Path]:
        """Convert PDF pages to images using PyMuPDF with enhanced quality"""
        images = []
        
        # Validate PDF file before processing
        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
            
        if pdf_path.stat().st_size == 0:
            raise ValueError(f"PDF file is empty: {pdf_path}")
        
        try:
            pdf_document = fitz.open(str(pdf_path))
            
            for page_num in range(pdf_document.page_count):
                page = pdf_document[page_num]
                
                # Calculate optimal matrix for rendering
                # Use higher resolution for better quality
                scale = self.dpi / 72.0
                mat = fitz.Matrix(scale, scale)
                
                # Get pixmap with antialiasing and better color depth
                if self.antialiasing:
                    pix = page.get_pixmap(
                        matrix=mat,
                        alpha=False,  # No transparency needed for PPT slides
                        colorspace=fitz.csRGB,  # Ensure RGB colorspace
                    )
                else:
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                
                # Convert to PIL Image for additional processing
                img_data = pix.pil_tobytes(format="PNG")
                img = Image.open(BytesIO(img_data))
                
                # Apply additional enhancements if needed
                if self.preserve_fonts:
                    # Ensure high quality when saving
                    save_kwargs = {
                        'quality': 95,
                        'optimize': True,
                        'dpi': (self.dpi, self.dpi)
                    }
                else:
                    save_kwargs = {'quality': 90, 'dpi': (self.dpi, self.dpi)}
                
                # Save image with enhanced quality
                image_path = output_dir / f"{prefix}_{page_num + 1:03d}.{image_format}"
                
                if image_format.lower() in ['jpg', 'jpeg']:
                    img = img.convert('RGB')  # Remove alpha channel for JPEG
                    img.save(str(image_path), format='JPEG', **save_kwargs)
                else:
                    img.save(str(image_path), format=image_format.upper(), **save_kwargs)
                
                images.append(image_path)
                logger.info(f"Saved slide {page_num + 1} as {image_path} with DPI: {self.dpi}")
                
            pdf_document.close()
            
        except Exception as e:
            logger.error(f"Error converting PDF to images: {e}")
            raise
            
        return images
    
    def optimize_for_chinese_text(self):
        """
        Optimize parser settings for Chinese text rendering
        """
        # Increase DPI for better Chinese character rendering
        if self.dpi < 300:
            self.dpi = 300
            logger.info(f"Increased DPI to {self.dpi} for better Chinese text rendering")
        
        # Enable font preservation and antialiasing
        self.preserve_fonts = True
        self.antialiasing = True
        
    def get_slide_metadata(self, ppt_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extract metadata from PPT/PPTX file
        
        Args:
            ppt_path: Path to PPT/PPTX file
            
        Returns:
            Dictionary containing slide metadata
        """
        ppt_path = Path(ppt_path)
        metadata = {
            'filename': ppt_path.name,
            'format': ppt_path.suffix.lower(),
            'slides': []
        }
        
        if ppt_path.suffix.lower() == '.pptx':
            try:
                prs = Presentation(str(ppt_path))
                metadata['slide_count'] = len(prs.slides)
                metadata['slide_width'] = prs.slide_width
                metadata['slide_height'] = prs.slide_height
                
                for i, slide in enumerate(prs.slides):
                    slide_info = {
                        'slide_number': i + 1,
                        'has_title': False,
                        'title': '',
                        'shape_count': len(slide.shapes)
                    }
                    
                    # Try to extract title
                    if slide.shapes.title:
                        slide_info['has_title'] = True
                        slide_info['title'] = slide.shapes.title.text
                    
                    metadata['slides'].append(slide_info)
                    
            except Exception as e:
                logger.warning(f"Could not extract metadata from PPTX: {e}")
        
        return metadata
    
    def parse_to_pil_images(
        self,
        ppt_path: Union[str, Path]
    ) -> List[Image.Image]:
        """
        Parse PPT/PPTX and return PIL Image objects instead of saving to disk
        
        Args:
            ppt_path: Path to PPT/PPTX file
            
        Returns:
            List of PIL Image objects
        """
        with tempfile.TemporaryDirectory() as tmp_dir:
            image_paths = self.parse(ppt_path, tmp_dir)
            
            images = []
            for path in image_paths:
                img = Image.open(path)
                images.append(img.copy())  # Copy to keep in memory
                
            return images